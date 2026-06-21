#!/usr/bin/env python3
"""
Создание современной PowerPoint презентации — Titanic Data Analysis
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_DIR = '/tmp/data-analysis-project/output/'
SLIDES_DIR = '/tmp/data-analysis-project/slides/'
os.makedirs(SLIDES_DIR, exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Цветовая схема (dark modern)
BG_DARK = RGBColor(18, 18, 32)
BG_LIGHT = RGBColor(245, 247, 250)
ACCENT = RGBColor(108, 92, 231)  # фиолетовый
ACCENT2 = RGBColor(0, 206, 201)  # бирюзовый
WHITE = RGBColor(255, 255, 255)
DARK_TEXT = RGBColor(33, 37, 41)
GRAY = RGBColor(160, 165, 175)
RED = RGBColor(255, 107, 107)
GREEN = RGBColor(46, 213, 115)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_accent_bar(slide, left=0.5, top=0.3, width=3, color=ACCENT):
    """Цветная полоска для заголовка"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(left), Inches(top), 
        Inches(width), Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_text_box(slide, left, top, width, height, text, size=20, bold=False, 
                 color=WHITE, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf

def add_bullet_text(slide, left, top, width, height, items, size=22, 
                    color=DARK_TEXT, spacing=12, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Эмодзи как буллиты
        if item.startswith('✓'):
            p.text = item
            p.font.size = Pt(size)
            p.font.color.rgb = GREEN
        elif item.startswith('✗'):
            p.text = item
            p.font.size = Pt(size)
            p.font.color.rgb = RED
        elif item.startswith('  •'):
            p.text = item
            p.font.size = Pt(size - 2)
            p.font.color.rgb = GRAY
        else:
            p.text = item
            p.font.size = Pt(size)
            p.font.color.rgb = color
        
        p.font.name = font_name
        p.space_after = Pt(spacing)
    return tf

def add_image_safe(slide, path, left, top, width, height, border=False):
    if os.path.exists(path):
        pic = slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
        if border:
            # Добавляем белую рамку
            from pptx.oxml.ns import qn
            spPr = pic._element.spPr
            ln = spPr.makeelement(qn('a:ln'), {'w': '6350', 'cap': 'flat'})
            solidFill = ln.makeelement(qn('a:solidFill'), {})
            srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': 'FFFFFF'})
            solidFill.append(srgbClr)
            ln.append(solidFill)
            spPr.append(ln)
        return True
    return False

def add_card(slide, left, top, width, height, color=RGBColor(255, 255, 255)):
    """Белая карточка с тенью"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(left), Inches(top), 
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Тень
    from pptx.oxml.ns import qn
    spPr = shape._element.spPr
    # Добавляем эффект тени
    effectDag = spPr.makeelement(qn('a:effectLst'), {})
    outerShdw = effectDag.makeelement(qn('a:outerShdw'), {
        'blurRad': '50800', 'dist': '12700', 'dir': '2700000',
        'algn': 'tl', 'rotWithShape': '0'
    })
    srgbClr = outerShdw.makeelement(qn('a:srgbClr'), {'val': '000000'})
    alpha = srgbClr.makeelement(qn('a:alpha'), {'val': '15000'})
    srgbClr.append(alpha)
    outerShdw.append(srgbClr)
    effectDag.append(outerShdw)
    spPr.append(effectDag)
    return shape

# ============ СЛАЙД 1: ТИТУЛ ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

# Акцентная линия сверху
add_accent_bar(slide, left=0, top=0, width=13.333, color=ACCENT)

add_text_box(slide, 1, 1.5, 11, 1, 'TITANIC', size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 1, 2.8, 11, 0.8, 'Data Analysis Project', size=36, color=ACCENT2, align=PP_ALIGN.CENTER)

# Разделитель
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.8), Inches(2.333), Inches(0.04))
div.fill.solid()
div.fill.fore_color.rgb = ACCENT
div.line.fill.background()

add_text_box(slide, 1, 4.2, 11, 0.5, 'Иода Алексей', size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.8, 11, 0.5, 'JA-группа | 2026', size=20, color=GRAY, align=PP_ALIGN.CENTER)

add_text_box(slide, 1, 5.8, 11, 0.5, 'Датасет: Titanic (крушение 1912) • 891 пассажир • 15 признаков', 
             size=16, color=GRAY, align=PP_ALIGN.CENTER)

# ============ СЛАЙД 2: О ДАННЫХ ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
add_accent_bar(slide, left=0.5, top=0.5, width=2, color=ACCENT)
add_text_box(slide, 0.5, 0.7, 12, 0.8, 'Объект исследования', size=36, bold=True, color=DARK_TEXT)

# Карточки
add_card(slide, 0.5, 1.8, 3.8, 2.5)
add_text_box(slide, 0.8, 2.0, 3.2, 0.5, '📊 Размер', size=24, bold=True, color=DARK_TEXT)
add_bullet_text(slide, 0.8, 2.6, 3.2, 1.5, ['891 строк', '15 столбцов', '6 числовых признаков'], size=20)

add_card(slide, 4.8, 1.8, 3.8, 2.5)
add_text_box(slide, 5.1, 2.0, 3.2, 0.5, '🎯 Цель', size=24, bold=True, color=DARK_TEXT)
add_bullet_text(slide, 5.1, 2.6, 3.2, 1.5, ['Предсказание выживания', 'Выявление факторов', 'Стат. анализ'], size=20)

add_card(slide, 9.1, 1.8, 3.8, 2.5)
add_text_box(slide, 9.4, 2.0, 3.2, 0.5, '🔍 Методы', size=24, bold=True, color=DARK_TEXT)
add_bullet_text(slide, 9.4, 2.6, 3.2, 1.5, ['PCA', 'KMeans', '10 гипотез'], size=20)

add_card(slide, 0.5, 4.8, 12.3, 2.2)
add_bullet_text(slide, 0.8, 5.0, 11.5, 1.8, [
    '✓ Источник: Seaborn библиотека (оригинал — Kaggle)',
    '✓ Целевая переменная: Survived (0 = погиб, 1 = выжил)',
    '✓ Признаки: класс, пол, возраст, стоимость билета, порт посадки, семья',
], size=22, color=DARK_TEXT)

# ============ СЛАЙД 3: ОПИСАТЕЛЬНАЯ СТАТИСТИКА ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
add_accent_bar(slide, left=0.5, top=0.5, width=2, color=ACCENT)
add_text_box(slide, 0.5, 0.7, 12, 0.8, 'Описательная статистика', size=36, bold=True, color=DARK_TEXT)

# Левая колонка — цифры
add_card(slide, 0.5, 1.8, 6, 5.2, RGBColor(255,255,255))
add_bullet_text(slide, 0.8, 2.0, 5.4, 3, [
    '📌 Средний возраст пассажира: 29.7 лет',
    '📌 Средняя стоимость билета: $32.20',
    '📌 Мужчин: 577 (64.8%) | Женщин: 314 (35.2%)',
    '📌 3 класс — 491 чел. (55.1%) | 1 класс — 216 (24.2%)',
    '📌 Выжило: 342 (38.4%) | Погибло: 549 (61.6%)',
    '',
    '⚠️ Пропуски: age (177), deck (688), embarked (2)',
], size=20, color=DARK_TEXT, spacing=8)

# Правая колонка — pie chart
add_image_safe(slide, f'{OUTPUT_DIR}6_pie.png', 7, 1.8, 5.5, 4.5)

# ============ СЛАЙД 4: ВИЗУАЛИЗАЦИЯ ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
add_accent_bar(slide, left=0.5, top=0.5, width=2, color=ACCENT)
add_text_box(slide, 0.5, 0.7, 12, 0.8, 'Визуальный анализ', size=36, bold=True, color=DARK_TEXT)

# Гистограммы
add_image_safe(slide, f'{OUTPUT_DIR}1_histograms.png', 0.5, 1.8, 6, 3)
# Box plots
add_image_safe(slide, f'{OUTPUT_DIR}2_boxplots.png', 7, 1.8, 5.8, 3)

# Нижний ряд — текст
add_card(slide, 0.5, 5.2, 5.5, 1.8, RGBColor(255,255,255))
add_bullet_text(slide, 0.8, 5.4, 5.0, 1.5, [
    '✓ Возраст: нормальное распределение (20-38 лет)',
    '✓ Fare: сильный перекос (большинство ≤ $50)',
], size=18, color=DARK_TEXT, spacing=4)

add_card(slide, 6.5, 5.2, 6.3, 1.8, RGBColor(255,255,255))
add_bullet_text(slide, 6.8, 5.4, 5.7, 1.5, [
    '✓ Box Plot: 13% выбросов по стоимости билета',
    '✓ Корреляция: fare +0.26, pclass -0.34 с survived',
], size=18, color=DARK_TEXT, spacing=4)


# ============ СЛАЙД 5: ГИПОТЕЗЫ 1/2 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
add_accent_bar(slide, left=0.5, top=0.5, width=2, color=ACCENT)
add_text_box(slide, 0.5, 0.7, 12, 0.8, 'Проверка гипотез (часть 1)', size=36, bold=True, color=DARK_TEXT)

add_card(slide, 0.5, 1.8, 12.3, 5.2, RGBColor(255,255,255))
add_bullet_text(slide, 0.8, 2.0, 11.5, 4.8, [
    '1️⃣ Средний возраст = 30 лет?',
    '   → One-Sample t-test, t=-0.55, p=0.58 → H₀ ПРИНЯТА (возраст ≈ 30)',
    '',
    '2️⃣ Возраст мужчин ≠ возраст женщин?',
    '   → Two-Sample t-test, t=2.50, p=0.013 → H₀ ОТКЛОНЕНА (мужчины старше)',
    '',
    '3️⃣ Количество сиблингов = количество родителей?',
    '   → Paired t-test, t=3.97, p=0.0001 → H₀ ОТКЛОНЕНА (разные величины)',
    '',
    '4️⃣ Средняя стоимость билета = $32?',
    '   → One-Sample z-test, z=0.12, p=0.90 → H₀ ПРИНЯТА (fare ≈ $32)',
    '',
    '5️⃣ Возраст различается по классам?',
    '   → ANOVA, F=57.44, p=0.0000 → H₀ ОТКЛОНЕНА (классы ≠ возраст)',
], size=20, color=DARK_TEXT, spacing=4)

# ============ СЛАЙД 6: ГИПОТЕЗЫ 2/2 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
add_accent_bar(slide, left=0.5, top=0.5, width=2, color=ACCENT)
add_text_box(slide, 0.5, 0.7, 12, 0.8, 'Проверка гипотез (часть 2)', size=36, bold=True, color=DARK_TEXT)

add_card(slide, 0.5, 1.8, 12.3, 5.2, RGBColor(255,255,255))
add_bullet_text(slide, 0.8, 2.0, 11.5, 4.8, [
    '6️⃣ Пол и выживание — зависимы?',
    '   → Chi² = 260.72, p=0.0000 → ДА (женщины выживают чаще) 🔴',
    '',
    '7️⃣ Класс и выживание — зависимы?',
    '   → Chi² = 102.89, p=0.0000 → ДА (1 класс выживает чаще) 🔴',
    '',
    '8️⃣ Стоимость билета выживших > погибших?',
    '   → Two-Sample t-test, t=7.94, p=0.0000 → ДА (выжившие платили больше) 🔴',
    '',
    '9️⃣ Порт посадки влияет на выживание?',
    '   → Chi² = 26.49, p=0.0000 → ДА (порт связан с классом) 🔴',
    '',
    '🔟 Средняя стоимость билета = $50?',
    '   → One-Sample t-test, t=-10.69, p=0.0000 → НЕТ (реально $32)',
], size=20, color=DARK_TEXT, spacing=4)

# ============ СЛАЙД 7: КОРРЕЛЯЦИЯ ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
add_accent_bar(slide, left=0.5, top=0.5, width=2, color=ACCENT)
add_text_box(slide, 0.5, 0.7, 12, 0.8, 'Корреляционный анализ', size=36, bold=True, color=DARK_TEXT)

add_image_safe(slide, f'{OUTPUT_DIR}7_correlation_heatmap.png', 0.5, 1.8, 7, 5.2)

add_card(slide, 8, 1.8, 4.8, 5.2, RGBColor(255,255,255))
add_bullet_text(slide, 8.3, 2.0, 4.2, 4.8, [
    'Корреляция с выживанием:',
    '',
    '✓ Fare: +0.26',
    '  (дороже билет → выше шанс)',
    '',
    '✗ Pclass: -0.34',
    '  (ниже класс → ниже шанс)',
    '',
    '○ Age: -0.08',
    '  (возраст почти не влияет)',
    '',
    'Вывод: социальный статус',
    'важнее возраста',
], size=20, color=DARK_TEXT, spacing=4)

# ============ СЛАЙД 8: PCA ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
add_accent_bar(slide, left=0.5, top=0.5, width=2, color=ACCENT)
add_text_box(slide, 0.5, 0.7, 12, 0.8, 'PCA — снижение размерности', size=36, bold=True, color=DARK_TEXT)

add_image_safe(slide, f'{OUTPUT_DIR}8_pca_variance.png', 0.5, 1.8, 6, 2.8)
add_image_safe(slide, f'{OUTPUT_DIR}9_pca_projection.png', 7, 1.8, 5.8, 2.8)

add_card(slide, 0.5, 5.0, 5.5, 2, RGBColor(255,255,255))
add_bullet_text(slide, 0.8, 5.2, 5.0, 1.6, [
    '📊 5 компонент = 90% информации',
    '📊 PC1+PC2 = 58.7% дисперсии',
], size=18, color=DARK_TEXT)

add_card(slide, 6.5, 5.0, 6.3, 2, RGBColor(255,255,255))
add_bullet_text(slide, 6.8, 5.2, 5.7, 1.5, [
    '🔵 Проекция: разделение по выживаемости',
    '🔵 PCA подтверждает структуру данных',
], size=18, color=DARK_TEXT)

# ============ СЛАЙД 9: KMEANS ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
add_accent_bar(slide, left=0.5, top=0.5, width=2, color=ACCENT)
add_text_box(slide, 0.5, 0.7, 12, 0.8, 'KMeans кластеризация', size=36, bold=True, color=DARK_TEXT)

add_image_safe(slide, f'{OUTPUT_DIR}10_kmeans_clusters.png', 0.5, 1.8, 7, 5)

add_card(slide, 8, 1.8, 4.8, 5.2, RGBColor(255,255,255))
add_bullet_text(slide, 8.3, 2.0, 4.2, 4.8, [
    'KMeans, k=2 на PCA-данных',
    '',
    'Кластер 0:',
    '   342 погибших + 144 выживших',
    '',
    'Кластер 1:',
    '   82 погибших + 146 выживших',
    '',
    'Социальное расслоение →',
    '2 группы пассажиров',
    '',
    'Кластеры коррелируют',
    'с выживаемостью ✅',
], size=20, color=DARK_TEXT, spacing=4)

# ============ СЛАЙД 10: CATPLOT ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
add_accent_bar(slide, left=0.5, top=0.5, width=2, color=ACCENT)
add_text_box(slide, 0.5, 0.7, 12, 0.8, 'Выживаемость: пол + класс', size=36, bold=True, color=DARK_TEXT)

add_image_safe(slide, f'{OUTPUT_DIR}5_catplot.png', 0.5, 1.8, 7.5, 5)

add_card(slide, 8.5, 1.8, 4.3, 5.2, RGBColor(255,255,255))
add_bullet_text(slide, 8.8, 2.0, 3.7, 4.5, [
    '🔴 1 класс:',
    '   ~95% женщин выжили',
    '',
    '🟡 2 класс:',
    '   ~90% женщин выжили',
    '',
    '🟢 3 класс:',
    '   ~50% женщин выжили',
    '',
    'Мужчины 3 класса:',
    '   <15% выжили',
], size=20, color=DARK_TEXT, spacing=6)

# ============ СЛАЙД 11: ВЫВОДЫ ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_accent_bar(slide, left=0, top=0, width=13.333, color=ACCENT)

add_text_box(slide, 0.5, 1, 12, 1, 'Ключевые выводы', size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_card(slide, 0.5, 2.5, 12.3, 4.5, RGBColor(30, 30, 60))
add_bullet_text(slide, 1, 2.8, 11, 3.8, [
    '✅ Проведена полная описательная статистика и визуализация данных',
    '✅ Выполнены 10 статистических тестов (t-test, ANOVA, χ²)',
    '✅ PCA — снижение размерности с 6 до 5 компонент (90% дисперсии)',
    '✅ KMeans — кластеризация выявила 2 социальные группы пассажиров',
    '',
    '🏆 Главное открытие: выживание на Титанике определялось',
    '   ПОЛОМ и СОЦИАЛЬНЫМ СТАТУСОМ, а не возрастом',
], size=22, color=WHITE, spacing=10)


# ============ СЛАЙД 12: ФИНАЛ ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_accent_bar(slide, left=0, top=0, width=13.333, color=ACCENT2)

add_text_box(slide, 1, 2, 11, 1.2, 'Спасибо за внимание!', size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

div2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.5), Inches(2.333), Inches(0.04))
div2.fill.solid()
div2.fill.fore_color.rgb = ACCENT2
div2.line.fill.background()

add_text_box(slide, 1, 4, 11, 0.8, 'Иода Алексей', size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.7, 11, 0.6, 'JA-группа', size=22, color=GRAY, align=PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.5, 11, 0.5, 'Вопросы?', size=24, color=ACCENT2, align=PP_ALIGN.CENTER)

prs.save(f'{SLIDES_DIR}Titanic_Data_Analysis_v2.pptx')
print(f"Презентация v2 сохранена: {SLIDES_DIR}Titanic_Data_Analysis_v2.pptx")
