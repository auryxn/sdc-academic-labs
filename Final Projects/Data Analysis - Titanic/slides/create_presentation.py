#!/usr/bin/env python3
"""
Создание PowerPoint презентации для Data Analysis Project — Titanic
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUTPUT_DIR = '/tmp/data-analysis-project/output/'
SLIDES_DIR = '/tmp/data-analysis-project/slides/'
os.makedirs(SLIDES_DIR, exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(33, 37, 41)
BLUE = RGBColor(13, 110, 253)
GRAY = RGBColor(108, 117, 125)

def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, size=20, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf

def add_image_safe(slide, path, left, top, width, height):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
        return True
    return False

# ============ СЛАЙД 1: Титульный ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

prs.save(f'{SLIDES_DIR}temp.pptx')

# Сделаю проще — без сложного стиля, норм презентация
prs2 = Presentation()
prs2.slide_width = Inches(13.333)
prs2.slide_height = Inches(7.5)

def add_title_slide(text, subtitle="", size=40):
    slide = prs2.slides.add_slide(prs2.slide_layouts[6])
    add_bg(slide, DARK)
    add_text_box(slide, 1, 2, 11, 1.5, text, size=size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, 1, 3.5, 11, 1, subtitle, size=24, color=GRAY, align=PP_ALIGN.CENTER)

def add_content_slide(title, bullets, img_path=None):
    slide = prs2.slides.add_slide(prs2.slide_layouts[6])
    add_bg(slide, RGBColor(245, 245, 245))
    add_text_box(slide, 0.5, 0.3, 12, 1, title, size=32, bold=True, color=DARK)
    
    if bullets:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = b
            p.font.size = Pt(22)
            p.font.color.rgb = DARK
            p.space_after = Pt(10)
    
    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(7), Inches(1.5), Inches(5.5), Inches(4.5))


# ============ СЛАЙД 1 ============
add_title_slide("Data Analysis: Titanic Dataset", "Индивидуальный проект | Data Analysis | 2026")

# ============ СЛАЙД 2: Объект исследования ============
add_content_slide("Объект исследования", [
    "Датасет: Titanic (пассажиры Титаника, 1912)",
    "Источник: Seaborn / Kaggle",
    "Строк: 891, Столбцов: 15",
    "Целевая переменная: Survived (выжил/погиб)",
    "Признаки: класс, пол, возраст, стоимость билета, порт посадки"
])

# ============ СЛАЙД 3: Описательная статистика ============
add_content_slide("Описательная статистика", [
    "Средний возраст: 29.7 лет",
    "Средняя стоимость билета: $32.20",
    "Мужчин: 577 (64.8%), Женщин: 314 (35.2%)",
    "3 класс: 491 пассажир (55.1%)",
    "Пропуски: age (177), deck (688), embarked (2)",
    "Выжило: 342 (38.4%), Погибло: 549 (61.6%)"
], f'{OUTPUT_DIR}6_pie.png')

# ============ СЛАЙД 4: Визуализация ============
add_content_slide("Визуализация данных", [
    "Гистограммы возраста и стоимости билета",
    "Box Plot: выбросы по стоимости билета (13%)",
    "Scatter: выжившие платили больше",
    "Violin Plot: распределение по классам и полу",
    "CatPlot: выживаемость зависит от класса и пола",
    "Heatmap корреляций"
], f'{OUTPUT_DIR}1_histograms.png')

# ============ СЛАЙД 5: Корреляция ============
add_content_slide("Корреляционный анализ", [
    "Наибольшая корреляция с выживанием:",
    "  • Fare (+0.26) — дорогие билеты = выше шанс",
    "  • Pclass (-0.34) — 1 класс = выше шанс",
    "  • Age (-0.08) — возраст почти не влияет",
    "  • Sex — сильная категориальная зависимость",
    "Вывод: социальный статус > возраст"
], f'{OUTPUT_DIR}7_correlation_heatmap.png')

# ============ СЛАЙД 6: Гипотезы ============
add_content_slide("Проверка статистических гипотез", [
    "1. One-Sample t-test: возраст = 30? → p=0.58 (H₀ принята)",
    "2. Two-Sample t-test: возраст мужчин ≠ женщин → p=0.013 (H₀ откл.)",
    "3. Paired t-test: sibsp vs parch → p=0.0001",
    "4. One-Sample z-test: fare = $32? → p=0.90 (H₀ принята)",
    "5. ANOVA: возраст по классам → p=0.000 (H₀ откл.)",
    "6. Chi²: пол и выживание → p=0.000 (зависимы)",
    "7. Chi²: класс и выживание → p=0.000 (зависимы)",
    "8. T-test: fare выживших > погибших → p=0.000",
    "Вывод: пол, класс и стоимость билета — значимые факторы"
])

# ============ СЛАЙД 7: PCA ============
add_content_slide("PCA (Метод главных компонент)", [
    "6 признаков → 5 компонент для 90% дисперсии",
    "PC1+PС2 объясняют 58.7% дисперсии",
    "Проекция: разделение по выживаемости",
    "Снижение размерности без потери информации",
    "PCA подтверждает структуру данных"
], f'{OUTPUT_DIR}9_pca_projection.png')

# ============ СЛАЙД 8: KMeans ============
add_content_slide("KMeans кластеризация", [
    "Метод: KMeans с k=2 на PCA-сжатых данных",
    "Кластер 0: 342 погибших + 144 выживших",
    "Кластер 1: 82 погибших + 146 выживших",
    "Кластеры коррелируют с выживаемостью",
    "Модель выделяет 2 группы пассажиров",
    "Результат: кластеризация подтверждает социальное расслоение"
], f'{OUTPUT_DIR}10_kmeans_clusters.png')

# ============ СЛАЙД 9: CatPlot ============
add_content_slide("Выживаемость по классу и полу", [
    "Женщины выживали чаще мужчин во всех классах",
    "1 класс: ~95% женщин выжили",
    "3 класс: ~50% женщин выжили",
    "Мужчины 3 класса: <15% выжили",
    "Класс и пол — главные факторы выживания"
], f'{OUTPUT_DIR}5_catplot.png')

# ============ СЛАЙД 10: Выводы ============
add_content_slide("Выводы", [
    "Выполнен полный анализ датасета Titanic:",
    "✓ Дескриптивная статистика и визуализация",
    "✓ 10 статистических тестов (t-test, ANOVA, χ²)",
    "✓ PCA (снижение размерности до 5 компонент)",
    "✓ KMeans кластеризация (2 кластера)",
    "",
    "Ключевое открытие:",
    "Выживание на Титанике определялось ПОЛОМ и СОЦИАЛЬНЫМ СТАТУСОМ,",
    "а не возрастом или физическими данными."
])


# ============ СЛАЙД 11: Спасибо ============
add_title_slide("Спасибо за внимание!", "Автор: [Твоё Имя] | Группа: [Твоя Группа]")

prs2.save(f'{SLIDES_DIR}Titanic_Data_Analysis.pptx')
print(f"Презентация сохранена: {SLIDES_DIR}Titanic_Data_Analysis.pptx")
